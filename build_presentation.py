import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define Colors
BG_COLOR = RGBColor(9, 9, 11)        # Zinc 950
TEXT_COLOR = RGBColor(244, 244, 245)  # Zinc 100
ACCENT_COLOR = RGBColor(16, 185, 129) # Emerald 500
SUBTEXT_COLOR = RGBColor(161, 161, 170) # Zinc 400

def set_background(slide, color=BG_COLOR):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, top=0.8, left=1.0, width=14.0, height=1.5, size=54, color=TEXT_COLOR):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    return txBox

def add_subtitle(slide, text, top=2.5, left=1.0, width=14.0, height=1.0, size=28, color=ACCENT_COLOR):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(size)
    p.font.color.rgb = color
    return txBox

def add_body(slide, text, top=3.5, left=1.0, width=14.0, height=4.5, size=22, color=SUBTEXT_COLOR, line_spacing=1.5):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(text.split('\n')):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.line_spacing = line_spacing
    return txBox

def add_image(slide, image_path, left, top, width=None, height=None):
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width) if width else None, height=Inches(height) if height else None)
    else:
        # Placeholder if image missing
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width if width else 4), Inches(height if height else 3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(50, 50, 50)
        p = shape.text_frame.paragraphs[0]
        p.text = "Image Placeholder"
        p.font.color.rgb = RGBColor(200, 200, 200)

blank_slide_layout = prs.slide_layouts[6]

# --- Slide 1: Title ---
slide = prs.slides.add_slide(blank_slide_layout)
set_background(slide)
add_title(slide, "杉蔓有限公司", top=3.0, size=72)
add_subtitle(slide, "SUSTAINABLE MAN CO., LTD.", top=4.5, size=32)
add_body(slide, "淨零時代，高效永續轉型策略夥伴\n\n品牌全書 & 公司介紹", top=5.5, size=24, color=TEXT_COLOR)
add_image(slide, "public/images/logo.png", left=12.5, top=0.5, width=3)
add_image(slide, "public/images/hero_2026.png", left=9.0, top=2.5, width=6.5)

# --- Slide 2: Origin ---
slide = prs.slides.add_slide(blank_slide_layout)
set_background(slide)
add_title(slide, "品牌起源：從嘉南平原的紅土地出發")
add_subtitle(slide, "「杉蔓的靈魂，誕生於烈日下的田間。」")
content = """
杉蔓有限公司並非始於冰冷的都市辦公室，而是起源於南台灣嘉南平原。

創辦團隊在投身「地方創生」的歲月中，親眼見證了氣候變遷對農業的衝擊：
老農在暴雨後看著倒伏作物的嘆息、優質產地因數位斷層而逐漸凋零。

我們深信：如果永續只是寫在報告書上的數字，而無法守護腳下這片土地，
那所有的商業繁榮都只是脆弱的海市蜃樓。
"""
add_body(slide, content.strip(), width=8.0)
add_image(slide, "public/images/about.png", left=9.5, top=2.5, width=5.5)

# --- Slide 3: Philosophy ---
slide = prs.slides.add_slide(blank_slide_layout)
set_background(slide)
add_title(slide, "品牌哲學：高聳與柔韌的結合")
add_subtitle(slide, "專業精準 與 無限連結")
content = """
「杉蔓」二字，代表了我們對永續事業的雙重承諾：

🌲「杉」 (Fir)
象徵杉木挺拔剛直、經冬不凋。
代表我們顧問服務的精準度、誠信與專業標準，
為企業建立永恆的基準。

🌿「蔓」 (Vine)
象徵藤蔓強韌的生命力與無限的連結。
致力於建構專業媒合平台，鏈結城市決策者與鄉村生產者，
形成共生共榮的永續生態圈。
"""
add_body(slide, content.strip(), width=8.0)
add_image(slide, "public/images/bg-ai.png", left=9.5, top=2.5, width=5.5)

# --- Slide 4: Vision & Mission ---
slide = prs.slides.add_slide(blank_slide_layout)
set_background(slide)
add_title(slide, "核心願景與使命")
content = """
🎯 品牌願景
做一個對地球有責任感、對下一代有溫度的永續人 (Sustainable Man)。

🚀 品牌使命
用溫暖的科技與創新的教育，解開數據的鎖，降低轉型門檻。

🤝 社會承諾 (Profit for Purpose)
承諾將年度盈餘的 15% 投入社會使命，
支持傳統農村轉型、偏鄉永續教育及人才培育。
"""
add_body(slide, content.strip(), top=2.5, width=14.0)

# --- Team Section: One slide per member ---
import requests
from io import BytesIO

team_members = [
    {
        "name": "曾偉庭 (Wei-Ting Tseng)",
        "role": "創辦人 暨 地方創生實踐家",
        "specialties": [
            "ISO14064-1 / ISO14067 主導查證員 (Lead Verifier)", "GRI 準則實務", "碳管理師 / 永續管理師", "IPAS 淨零碳規劃管理師",
            "杉蔓有限公司 負責人", "自由落腳商模有限公司 永續管理總監", "台灣碳盤查工會 副秘書長", "新光匯盈有限公司 經理/永續講師",
            "救國團永續發展研究院 執行顧問", "商業研究發展院 講師", "SDGs 永續教學講師", "永續桌遊培訓師", "農業永續轉型輔導講師", "地方創生實作坊 講師"
        ],
        "image": "https://drive.google.com/thumbnail?id=16mOnqgwo2Y0Xa2qEdQZwy2Dp084qiWji&sz=w500",
        "local_image": "public/images/team/tseng.jpg"
    },
    {
        "name": "李崇毅",
        "role": "資深策略顧問",
        "specialties": [
            "國立成功大學鳳凰新創平台新創導師", "台灣碳盤查產業工會名譽理事長", "奔馳科技股份有限公司執行長", 
            "ESG永續發展協會綠建築委員會主任委員", "國立成功大學品質與創新研究中心總顧問", "中華民國產業聯合發展總會秘書處執行長",
            "冠頡管理顧問有限公司執行長", "宏大文教事業股份有限公司董事", "煊程科技股份有限公司顧問暨新創團隊導師",
            "大員生醫科技股份有限公司新創團隊導師", "中華談判研究發展協會組織發展組副主任委員", "加鑫管理顧問有限公司執行長", "全珈科技股份有限公司總顧問"
        ],
        "image": None,
        "local_image": "public/images/team/li.jpg"
    },
    {
        "name": "鄭安姵",
        "role": "資深永續輔導顧問",
        "specialties": [
            "曾任國際驗證主任稽核員", "經濟部iPAS淨零碳規劃管理師", "SRS Social Reporting Standard level 1", 
            "ISO 14001 :2015環境管理系統主任稽核員", "ISO 14064-1: 2018溫室氣體盤查", "ISO 14067: 2018 碳足跡主任查證員",
            "ISO 14068-1: 2023 碳中和主任查證員", "ISO 20400: 2017 永續採購", "ISO 32210 永續金融-企業永續風險與機會系統性管理框架",
            "ISO 5001:2018能源管理系統主任稽核員", "GHG Protocol 會計法則 UCS", "SBTi 標準指南及工具培訓 UCS", "TTQS人才發展品質教育管理系統"
        ],
        "image": "https://drive.google.com/thumbnail?id=1EsDULLlQL2uHVyvT_TiYYl5AwttrxLRV&sz=w500",
        "local_image": "public/images/team/zheng.jpg"
    },
    {
        "name": "周家豪",
        "role": "資深永續碳權顧問",
        "specialties": [
            "益珂環能股份有限公司總經理", "COP29 UNFCCC 特邀講者", "Strathclyde U,UK 行銷博士", "Loughborough U,UK 行銷碩士",
            "2010年參與台灣第一個太陽能整廠輸出紀錄", "2012建置台灣第一個MW及地面型發電廠", "企業永續競爭力分析", 
            "商業模式創新設計", "全球行銷活動規劃專業執行", "再生能源暨綠能專業輔導顧問"
        ],
        "image": "https://drive.google.com/thumbnail?id=1D0g4Dwb0t4WoEPIme65ArsfDPWhP9OsI&sz=w500",
        "local_image": "public/images/team/zhou.jpg"
    },
    {
        "name": "詹家和",
        "role": "資深永續輔導顧問",
        "specialties": [
            "易德福國際有限公司總經理", "國立勤益科技大學ESG永續發展中心副執行長", "中華創新發展協會 ESG管理期刊主編",
            "中華永續科技創新發展協會祕書長", "專利：用於管理交易標的碳足跡的區塊鏈會計系統(I872336)", 
            "專利：碳價計算管理系統(M640876)", "專利：碳決策管理系統(M652062)", "台灣國際商貿協進會監事", 
            "中華永續商道聯合會創會長", "薩摩亞商萬達環球股份有限公司總經理"
        ],
        "image": None,
        "local_image": "public/images/team/zhan.jpg"
    },
    {
        "name": "陳淑真",
        "role": "資深永續輔導顧問",
        "specialties": [
            "始亦真永續有限公司 執行長", "國際ESG永續發展協會 理事", "濱川企業 永續發展管理代表", "義守大學溫盤諮詢及授課顧問",
            "經濟部產業低碳化輔導計畫-顧問(I872336)", "經濟部軟體協會國際行銷-顧問師", "2023經濟部工業局碳盤查加值應用計畫-共同計畫主持人",
            "ISO 14064-1 組織溫室氣體主任查證員", "ISO 14067 產品碳足跡主任查證員", "勞動部大小人提-企業內聘講師"
        ],
        "image": None,
        "local_image": "public/images/team/chen.jpg"
    }
]

# Create a section intro slide for Team
slide = prs.slides.add_slide(blank_slide_layout)
set_background(slide)
add_title(slide, "精英顧問團隊", top=3.5, size=60, left=1.0)
add_subtitle(slide, "匯聚跨領域的頂尖專家，結合學術深度與實務經驗", top=5.0, size=30, left=1.0)
add_image(slide, "public/images/logo.png", left=12.5, top=0.5, width=3)

# Create a slide for each member
for member in team_members:
    slide = prs.slides.add_slide(blank_slide_layout)
    set_background(slide)
    add_title(slide, member["name"], top=1.0, size=50)
    add_subtitle(slide, member["role"], top=2.2, size=30, color=ACCENT_COLOR)
    
    # Format specialties
    specs_text = "\n".join([f"• {s}" for s in member["specialties"]])
    # Use smaller font size and line spacing for longer lists
    font_size = 18 if len(member["specialties"]) > 8 else 22
    line_sp = 1.1 if len(member["specialties"]) > 8 else 1.5
    
    # Split to 2 columns if more than 10 items to make it look better
    if len(member["specialties"]) > 8:
        mid = len(member["specialties"]) // 2 + (len(member["specialties"]) % 2 > 0)
        col1_text = "\n".join([f"• {s}" for s in member["specialties"][:mid]])
        col2_text = "\n".join([f"• {s}" for s in member["specialties"][mid:]])
        add_body(slide, col1_text, top=3.5, left=1.0, width=5.0, size=16, line_spacing=1.2)
        add_body(slide, col2_text, top=3.5, left=6.5, width=5.0, size=16, line_spacing=1.2)
    else:
        add_body(slide, specs_text, top=3.5, left=1.0, width=8.0, size=font_size, line_spacing=line_sp)
    
    # Handle Image
    img_path = member["local_image"]
    if not os.path.exists(img_path) and member.get("image"):
        try:
            res = requests.get(member["image"])
            if res.status_code == 200:
                os.makedirs(os.path.dirname(img_path), exist_ok=True)
                with open(img_path, 'wb') as f:
                    f.write(res.content)
        except Exception as e:
            print(f"Error downloading image for {member['name']}: {e}")
            
    if os.path.exists(img_path):
        add_image(slide, img_path, left=12.0, top=2.0, width=3.5)
    else:
        # Placeholder
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.0), Inches(2.0), Inches(3.5), Inches(4.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(50, 50, 50)
        p = shape.text_frame.paragraphs[0]
        p.text = "照片準備中"
        p.font.color.rgb = RGBColor(200, 200, 200)


# --- Slide 6: Core Business 1 ---
slide = prs.slides.add_slide(blank_slide_layout)
set_background(slide)
add_title(slide, "核心業務 (1/2)")
content = """
💡 1. MS-RDK 空調專利節能解決方案 (旗艦產品)
• 獨家技術：日本專利陶瓷遠紅外線技術 + AI 精準模擬。
• 實質成效：一分鐘免工地安裝，實現 15%~25% 的 Scope 2 實質減碳。

📊 2. ISO 國際標準認證輔導
• 範疇：ISO 14064、ISO 14067、ISO 14068。
• 特色：一站式數據蒐集至報告書產出，確保接軌國際綠色供應鏈。
"""
add_body(slide, content.strip(), top=2.5, width=8.0)
add_image(slide, "public/images/msrdk_solution.png", left=9.5, top=2.5, width=5.5)

# --- Slide 7: Core Business 2 ---
slide = prs.slides.add_slide(blank_slide_layout)
set_background(slide)
add_title(slide, "核心業務 (2/2)")
content = """
🎲 3. ESG 策略桌遊培訓 (教育創新)
• 將複雜政策轉化為模擬決策桌遊，降低員工學習門檻，培育種子教師。

🤝 4. 專業媒合平台與顧問
• 協助企業尋找永續供應鏈夥伴，對接數位綠色雙軸轉型資源。

♻️ 5. 綠色策展與循環印刷
• 使用 rPET 再生布料與無毒印刷，解決傳統展場的浪費問題。
"""
add_body(slide, content.strip(), top=2.5, width=8.0)
add_image(slide, "public/images/services/esg_game.png", left=9.5, top=2.5, width=5.5)

# --- Slide 8: Advantages ---
slide = prs.slides.add_slide(blank_slide_layout)
set_background(slide)
add_title(slide, "我們的競爭優勢")
content = """
1️⃣ 實效導向
不只提供諮詢報告，更提供專利節能硬體 (MS-RDK)，讓減碳看得到。

2️⃣ 教育賦能
透過桌遊與課程讓員工從「被動合規」轉向「主動創新」。

3️⃣ 跨域整合
具備從農業創生、產學研究到國際碳權的完整生態圈資源。

4️⃣ 合規權威
顧問團隊皆具備國際主導查證員證照，確保輔導流程的嚴謹性。
"""
add_body(slide, content.strip(), top=2.5, width=14.0)

# --- Slide 9: Impact ---
slide = prs.slides.add_slide(blank_slide_layout)
set_background(slide)
add_title(slide, "社會影響力與行動紀錄")
add_subtitle(slide, "杉蔓的腳印遍佈全台，並積極參與國際對話")
content = """
• 學術培育：受邀於政大、文化大學產碩專班講授永續課程。
• 政策倡議：於林業及自然保育署舉辦 TNFD 講座。
• 國際接軌：遠赴日本九州進行企業永續商談，引進跨國技術。
• 普惠教育：舉辦多場「桌遊種子教師班」，讓永續知識遍地開花。
"""
add_body(slide, content.strip(), width=9.0)
add_image(slide, "public/images/mockup-esg.png", left=10.0, top=3.0, width=5.0)

# --- Slide 10: Ending ---
slide = prs.slides.add_slide(blank_slide_layout)
set_background(slide)
add_title(slide, "杉蔓，為您紮根專業，蔓生永續。", top=4.0, size=50)
add_subtitle(slide, "Let's thrive together.", top=5.5, size=30)
add_image(slide, "public/images/logo.png", left=7.0, top=1.5, width=2)

prs.save("杉蔓有限公司_公司介紹簡報_v2.pptx")
print("Presentation saved successfully!")
